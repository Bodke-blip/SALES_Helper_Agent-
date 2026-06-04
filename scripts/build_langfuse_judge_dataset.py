from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from time import sleep
from typing import Any

import requests
from openpyxl import Workbook
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

from ingestion import normalize_drive_document_name, normalize_match_key
from upload_to_qdrant import (
    BENEFIT_LABELS,
    SOLUTION_LABELS,
    TOOLS_LABELS,
    extract_labeled_value,
    split_labeled_list,
)


DEFAULT_ZIP = Path("/Users/atharvabodke/Downloads/drive-download-20260604T050253Z-3-001.zip")
DEFAULT_OUTPUT = Path("data/langfuse_judge_eval_dataset.xlsx")
DEFAULT_AUDIT = Path("data/langfuse_judge_eval_dataset_audit.json")
DEFAULT_API_URL = "http://127.0.0.1:8000/query"
EXCEL_ILLEGAL_CHARACTERS = re.compile(r"[\x00-\x08\x0b-\x0c\x0e-\x1f]")
EXCEL_MAX_CELL_CHARS = 32767


@dataclass(frozen=True)
class SourceDeck:
    original_name: str
    normalized_name: str
    match_key: str
    extracted_path: Path
    text: str
    tools: list[str]
    solution: str
    benefits: list[str]


def clean_text(value: str) -> str:
    value = value.replace("\x00", "")
    value = EXCEL_ILLEGAL_CHARACTERS.sub(" ", value)
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def safe_sheet_text(shape: Any) -> str:
    if hasattr(shape, "text") and shape.text:
        return str(shape.text).strip()
    return ""


def extract_ppt_text(path: Path) -> str:
    presentation = Presentation(path)
    parts: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = []

        for shape in slide.shapes:
            text = safe_sheet_text(shape)
            if text:
                slide_parts.append(text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))

        if slide_parts:
            parts.append(f"Slide {slide_number}\n" + "\n".join(slide_parts))

    return clean_text("\n\n".join(parts))


def choose_preferred_file(candidates: list[zipfile.ZipInfo]) -> zipfile.ZipInfo:
    return sorted(
        candidates,
        key=lambda item: (
            item.filename != normalize_drive_document_name(item.filename),
            item.filename.lower().count("copy of"),
            item.filename.lower(),
            item.CRC,
        ),
    )[0]


def extract_deduped_decks(zip_path: Path, extract_dir: Path) -> tuple[list[tuple[zipfile.ZipInfo, Path]], dict[str, Any]]:
    extract_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(zip_path) as archive:
        ppt_infos = [
            info
            for info in archive.infolist()
            if not info.is_dir() and info.filename.lower().endswith(".pptx")
        ]
        grouped: dict[str, list[zipfile.ZipInfo]] = {}

        for info in ppt_infos:
            grouped.setdefault(normalize_match_key(info.filename), []).append(info)

        selected = []
        skipped = []

        for match_key, infos in sorted(grouped.items()):
            selected_info = choose_preferred_file(infos)
            normalized_name = normalize_drive_document_name(selected_info.filename)
            output_path = extract_dir / normalized_name
            output_path.parent.mkdir(parents=True, exist_ok=True)

            with archive.open(selected_info) as source:
                output_path.write_bytes(source.read())

            selected.append((selected_info, output_path))

            for info in infos:
                if info is selected_info:
                    continue
                skipped.append(
                    {
                        "match_key": match_key,
                        "skipped_name": info.filename,
                        "kept_name": selected_info.filename,
                    }
                )

    audit = {
        "zip_path": str(zip_path),
        "selected_count": len(selected),
        "skipped_duplicate_count": len(skipped),
        "skipped_duplicates": skipped,
    }
    return selected, audit


def build_source_decks(zip_path: Path, work_dir: Path) -> tuple[list[SourceDeck], dict[str, Any]]:
    selected_files, audit = extract_deduped_decks(zip_path, work_dir)
    decks = []

    for info, path in selected_files:
        text = extract_ppt_text(path)
        lines = text.splitlines()
        tools = split_labeled_list(extract_labeled_value(lines, TOOLS_LABELS, max_lines=8))
        solution = extract_labeled_value(lines, SOLUTION_LABELS, max_lines=10)
        benefits = split_labeled_list(extract_labeled_value(lines, BENEFIT_LABELS, max_lines=10))

        if not text:
            continue

        decks.append(
            SourceDeck(
                original_name=info.filename,
                normalized_name=normalize_drive_document_name(info.filename),
                match_key=normalize_match_key(info.filename),
                extracted_path=path,
                text=text,
                tools=tools,
                solution=solution,
                benefits=benefits,
            )
        )

    audit["usable_deck_count"] = len(decks)
    return decks, audit


def infer_title(deck: SourceDeck) -> str:
    return deck.normalized_name.removesuffix(".pptx").strip()


def compact_list(items: list[str], fallback: str = "Not specified in extracted source text.") -> str:
    cleaned = [item for item in items if item]
    return ", ".join(cleaned[:8]) if cleaned else fallback


def ground_truth_overview(deck: SourceDeck) -> str:
    excerpt = clean_text(deck.text[:1200])
    return (
        f"Source document: {deck.normalized_name}\n"
        f"Expected answer should discuss the use case from this document only. "
        f"Solution: {deck.solution or 'Not explicitly labeled in extracted text.'}\n"
        f"Tools: {compact_list(deck.tools)}\n"
        f"Benefits: {compact_list(deck.benefits)}\n"
        f"Evidence excerpt: {excerpt}"
    )


def ground_truth_tools(deck: SourceDeck) -> str:
    return (
        f"Source document: {deck.normalized_name}\n"
        f"Expected answer should identify tools/technologies if present and avoid inventing missing tools. "
        f"Tools in extracted source text: {compact_list(deck.tools)}\n"
        f"Relevant evidence: {clean_text(deck.text[:900])}"
    )


def ground_truth_benefits(deck: SourceDeck) -> str:
    return (
        f"Source document: {deck.normalized_name}\n"
        f"Expected answer should summarize the stated business benefits/outcomes and avoid unsupported metrics. "
        f"Benefits in extracted source text: {compact_list(deck.benefits)}\n"
        f"Relevant evidence: {clean_text(deck.text[:900])}"
    )


def build_eval_rows(decks: list[SourceDeck], target_count: int) -> list[dict[str, str]]:
    rows = []

    for deck in decks:
        title = infer_title(deck)
        rows.append(
            {
                "query": f"Explain the {title} use case with customer, problem, solution, tools, benefits, and source grounding.",
                "ground_truth": ground_truth_overview(deck),
                "source_document": deck.normalized_name,
            }
        )

        if len(rows) >= target_count:
            return rows

        rows.append(
            {
                "query": f"What tools or technologies are used in the {title} use case?",
                "ground_truth": ground_truth_tools(deck),
                "source_document": deck.normalized_name,
            }
        )

        if len(rows) >= target_count:
            return rows

        rows.append(
            {
                "query": f"What benefits or outcomes are described for the {title} use case?",
                "ground_truth": ground_truth_benefits(deck),
                "source_document": deck.normalized_name,
            }
        )

        if len(rows) >= target_count:
            return rows

    return rows


def call_agent(api_url: str, query: str, timeout: int) -> str:
    response = requests.post(
        api_url,
        json={
            "query": query,
            "force_db_search": True,
            "use_gemini_llm": True,
            "use_local_llm": True,
            "verbose": False,
            "answer_style": "detailed",
        },
        timeout=timeout,
    )
    response.raise_for_status()
    payload = response.json()
    return str(payload.get("answer") or payload.get("message") or "")


def add_excel_rows(output_path: Path, rows: list[dict[str, str]]) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Langfuse Judge Dataset"
    sheet.append(["query", "ground_truth", "agent_reply"])

    for row in rows:
        sheet.append(
            [
                excel_safe_text(row["query"]),
                excel_safe_text(row["ground_truth"]),
                excel_safe_text(row["agent_reply"]),
            ]
        )

    sheet.column_dimensions["A"].width = 70
    sheet.column_dimensions["B"].width = 90
    sheet.column_dimensions["C"].width = 90
    workbook.save(output_path)


def excel_safe_text(value: str) -> str:
    return clean_text(str(value))[:EXCEL_MAX_CELL_CHARS]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a Langfuse LLM-as-judge eval dataset from Drive export zip.")
    parser.add_argument("--zip", type=Path, default=DEFAULT_ZIP)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--audit", type=Path, default=DEFAULT_AUDIT)
    parser.add_argument("--api-url", default=DEFAULT_API_URL)
    parser.add_argument("--count", type=int, default=75)
    parser.add_argument("--timeout", type=int, default=120)
    parser.add_argument("--sleep", type=float, default=0.1)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    work_dir = Path("data/langfuse_judge_eval_sources")
    decks, audit = build_source_decks(args.zip, work_dir)
    rows = build_eval_rows(decks, args.count)

    for index, row in enumerate(rows, start=1):
        print(f"[{index}/{len(rows)}] {row['query']}")
        try:
            row["agent_reply"] = call_agent(args.api_url, row["query"], args.timeout)
        except Exception as error:
            row["agent_reply"] = f"ERROR: {error}"
        sleep(args.sleep)

    add_excel_rows(args.output, rows)

    audit.update(
        {
            "requested_eval_rows": args.count,
            "written_eval_rows": len(rows),
            "output_path": str(args.output),
        }
    )
    args.audit.parent.mkdir(parents=True, exist_ok=True)
    args.audit.write_text(json.dumps(audit, indent=2), encoding="utf-8")
    print(f"Wrote Excel dataset: {args.output}")
    print(f"Wrote audit: {args.audit}")


if __name__ == "__main__":
    main()
